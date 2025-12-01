"""
PowerPoint generator for executive summary and full findings decks.

Creates professional presentations from synthesized deliverables.
"""

import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..config import OUTPUT_DIR, DELIVERABLES
from ..models import SynthesisOutput, ResearchOutput, CompanyInput


# Color scheme
COLORS = {
    "dark_blue": RGBColor(31, 78, 121),      # Primary headers
    "light_blue": RGBColor(68, 114, 196),    # Secondary elements
    "gray": RGBColor(89, 89, 89),            # Subtitles
    "dark_gray": RGBColor(64, 64, 64),       # Body text
    "white": RGBColor(255, 255, 255),        # Light text
    "light_bg": RGBColor(242, 242, 242),     # Light backgrounds
    "accent": RGBColor(0, 176, 80),          # Success/highlights
    "warning": RGBColor(255, 192, 0),        # Warnings
}


class PowerPointGenerator:
    """
    Generates PowerPoint presentations from deliverables.

    Responsibilities:
    - Create executive summary deck (10-15 slides)
    - Create full findings presentation (30-40 slides)
    - Apply consistent styling and branding
    - Include mermaid diagram images
    """

    def __init__(self, output_dir: Optional[Path] = None):
        """
        Initialize the PowerPoint generator.

        Args:
            output_dir: Base output directory.
        """
        self.output_dir = output_dir or OUTPUT_DIR

    def generate_executive_summary(
        self,
        company_slug: str,
        company_input: CompanyInput,
        research: ResearchOutput,
        synthesis: SynthesisOutput,
        mermaid_images: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate executive summary deck (10-15 slides).

        Slides:
        1. Title slide
        2. Executive summary / Key findings
        3. Current state overview
        4. AI maturity assessment
        5. Key pain points
        6. Top 5 quick wins
        7. Strategic roadmap overview
        8. Expected ROI
        9. Recommended next steps
        10. Appendix / Contact

        Args:
            company_slug: URL-safe company name.
            company_input: Original company input.
            research: Research output.
            synthesis: Synthesis output.
            mermaid_images: Dict of diagram name to image path.

        Returns:
            Path to generated PPTX file.
        """
        prs = Presentation()

        # Set slide dimensions (widescreen 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Build slides
        self._add_title_slide(
            prs,
            company_input.name,
            "AI Strategy & Implementation Roadmap",
            f"Executive Summary | {datetime.now().strftime('%B %Y')}",
        )

        self._add_executive_summary_slide(prs, company_input, research, synthesis)
        self._add_current_state_slide(prs, synthesis, mermaid_images)
        self._add_maturity_slide(prs, synthesis)
        self._add_pain_points_slide(prs, synthesis)
        self._add_quick_wins_slide(prs, synthesis)
        self._add_roadmap_overview_slide(prs, synthesis)
        self._add_roi_slide(prs, synthesis)
        self._add_next_steps_slide(prs, company_input)
        self._add_contact_slide(prs)

        # Save presentation
        output_path = self._get_output_path(company_slug, "executive_summary.pptx")
        prs.save(output_path)

        return str(output_path)

    def generate_full_findings(
        self,
        company_slug: str,
        company_input: CompanyInput,
        research: ResearchOutput,
        synthesis: SynthesisOutput,
        mermaid_images: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate full findings presentation (30-40 slides).

        Comprehensive presentation covering all deliverables.

        Args:
            company_slug: URL-safe company name.
            company_input: Original company input.
            research: Research output.
            synthesis: Synthesis output.
            mermaid_images: Dict of diagram name to image path.

        Returns:
            Path to generated PPTX file.
        """
        prs = Presentation()

        # Set slide dimensions
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(
            prs,
            company_input.name,
            "AI Strategy Assessment",
            f"Full Findings & Recommendations | {datetime.now().strftime('%B %Y')}",
        )

        # Agenda
        self._add_agenda_slide(prs)

        # Section: Company Overview & Research
        self._add_section_divider(prs, "Company Overview & Market Context")
        self._add_company_profile_slide(prs, research)
        self._add_industry_context_slide(prs, research)
        self._add_competitive_landscape_slide(prs, research)

        # Section: Current State Assessment
        self._add_section_divider(prs, "Current State Assessment")
        self._add_tech_inventory_slides(prs, synthesis)
        self._add_pain_points_detailed_slide(prs, synthesis)
        self._add_current_state_diagram_slide(prs, mermaid_images)

        # Section: AI Maturity & Readiness
        self._add_section_divider(prs, "AI Maturity & Readiness")
        self._add_maturity_detailed_slides(prs, synthesis)

        # Section: Use Cases & Opportunities
        self._add_section_divider(prs, "AI Use Cases & Opportunities")
        self._add_use_case_library_slides(prs, synthesis)

        # Section: Strategy & Roadmap
        self._add_section_divider(prs, "Strategic Roadmap")
        self._add_quick_wins_detailed_slides(prs, synthesis)
        self._add_roadmap_detailed_slides(prs, synthesis)
        self._add_future_state_diagram_slide(prs, mermaid_images)

        # Section: Financial Analysis
        self._add_section_divider(prs, "Financial Analysis & ROI")
        self._add_roi_detailed_slides(prs, synthesis)
        self._add_vendor_comparison_slides(prs, synthesis)

        # Section: Governance & Operations
        self._add_section_divider(prs, "Governance & Operations")
        self._add_governance_slides(prs, synthesis)

        # Section: Change Management
        self._add_section_divider(prs, "Change Management & Training")
        self._add_change_management_slides(prs, synthesis)

        # Closing
        self._add_summary_slide(prs, synthesis)
        self._add_next_steps_slide(prs, company_input)
        self._add_appendix_slide(prs)
        self._add_contact_slide(prs)

        # Save presentation
        output_path = self._get_output_path(company_slug, "full_findings.pptx")
        prs.save(output_path)

        return str(output_path)

    def _get_output_path(self, company_slug: str, filename: str) -> Path:
        """Get output path for a presentation."""
        output_dir = self.output_dir / company_slug / "presentations"
        output_dir.mkdir(parents=True, exist_ok=True)
        return output_dir / filename

    # =========================================================================
    # Slide building methods
    # =========================================================================

    def _add_title_slide(
        self,
        prs: Presentation,
        company_name: str,
        title: str,
        subtitle: str,
    ) -> None:
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{company_name}"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS["dark_blue"]
        p.alignment = PP_ALIGN.CENTER

        # Add main title
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12.333), Inches(1)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = COLORS["dark_gray"]
        p.alignment = PP_ALIGN.CENTER

        # Add date/subtitle
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(12.333), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["gray"]
        p.alignment = PP_ALIGN.CENTER

    def _add_section_divider(self, prs: Presentation, title: str) -> None:
        """Add a section divider slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add colored background shape
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["dark_blue"]
        shape.line.fill.background()

        # Add section title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

    def _add_slide_with_title(
        self,
        prs: Presentation,
        title: str,
        subtitle: Optional[str] = None,
    ):
        """Add a content slide with title and return the slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS["dark_blue"]

        # Subtitle if provided
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS["gray"]

        return slide

    def _add_bullets_to_slide(
        self,
        slide,
        bullets: List[str],
        left: float = 0.5,
        top: float = 1.5,
        width: float = 12.333,
        height: float = 5.5,
        font_size: int = 18,
    ) -> None:
        """Add bullet points to a slide."""
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {text}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS["dark_gray"]
            p.space_after = Pt(12)

    def _add_table_to_slide(
        self,
        slide,
        headers: List[str],
        rows: List[List[str]],
        left: float = 0.5,
        top: float = 1.5,
        width: float = 12.333,
        height: float = 4,
    ) -> None:
        """Add a table to a slide."""
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        ).table

        # Set column widths
        col_width = width / num_cols
        for i in range(num_cols):
            table.columns[i].width = Inches(col_width)

        # Style header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["dark_blue"]

            # Style text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = COLORS["white"]
                paragraph.alignment = PP_ALIGN.CENTER

        # Add data rows
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)

                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS["light_bg"]

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.font.color.rgb = COLORS["dark_gray"]

    def _extract_content_section(
        self,
        synthesis: SynthesisOutput,
        deliverable_id: str,
        section_name: Optional[str] = None,
    ) -> str:
        """Extract content from a deliverable."""
        if deliverable_id not in synthesis.deliverables:
            return ""

        content = synthesis.deliverables[deliverable_id].content
        if not content:
            return ""

        if section_name:
            # Extract specific section
            pattern = rf'##\s*{re.escape(section_name)}.*?\n(.*?)(?=\n##|\Z)'
            match = re.search(pattern, content, re.DOTALL | re.IGNORECASE)
            if match:
                return match.group(1).strip()

        return content

    def _extract_bullets_from_content(
        self,
        content: str,
        max_bullets: int = 6,
    ) -> List[str]:
        """Extract bullet points from markdown content."""
        bullet_pattern = re.compile(r'^\s*[-*+]\s+(.+)$', re.MULTILINE)
        matches = bullet_pattern.findall(content)
        return [m.strip() for m in matches[:max_bullets]]

    def _extract_table_from_content(
        self,
        content: str,
    ) -> Optional[Dict[str, Any]]:
        """Extract first table from markdown content."""
        lines = content.split("\n")
        table_lines = []
        in_table = False

        for line in lines:
            line = line.strip()
            if line.startswith("|") and line.endswith("|"):
                in_table = True
                table_lines.append(line)
            elif in_table and not line.startswith("|"):
                break

        if len(table_lines) >= 3:
            # Parse headers
            headers = [c.strip() for c in table_lines[0].strip("|").split("|")]
            # Parse rows (skip separator)
            rows = []
            for line in table_lines[2:]:
                cells = [c.strip() for c in line.strip("|").split("|")]
                if len(cells) == len(headers):
                    rows.append(cells)
            return {"headers": headers, "rows": rows}

        return None

    # =========================================================================
    # Executive Summary specific slides
    # =========================================================================

    def _add_executive_summary_slide(
        self,
        prs: Presentation,
        company_input: CompanyInput,
        research: ResearchOutput,
        synthesis: SynthesisOutput,
    ) -> None:
        """Add executive summary slide."""
        slide = self._add_slide_with_title(prs, "Executive Summary")

        key_findings = [
            f"Industry: {research.industry.primary_industry}" if research.industry.primary_industry else "Technology-focused organization",
            f"AI Maturity: {research.tech_landscape.ai_maturity_estimate}" if research.tech_landscape.ai_maturity_estimate else "Early exploration phase",
            f"Key opportunity areas identified across {len(research.industry.opportunities)} industry trends",
            "Quick wins identified for immediate implementation",
            "Comprehensive roadmap spanning 30-360 days",
            "ROI potential with quantified cost-benefit analysis",
        ]

        self._add_bullets_to_slide(slide, key_findings)

    def _add_current_state_slide(
        self,
        prs: Presentation,
        synthesis: SynthesisOutput,
        mermaid_images: Optional[Dict[str, str]],
    ) -> None:
        """Add current state overview slide."""
        slide = self._add_slide_with_title(prs, "Current State Overview")

        # Try to add diagram image
        if mermaid_images and "current_state" in mermaid_images:
            image_path = Path(mermaid_images["current_state"])
            if image_path.exists() and image_path.suffix == ".png":
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(1), Inches(1.5),
                    width=Inches(11),
                )
                return

        # Fallback to text bullets
        content = self._extract_content_section(synthesis, "01_tech_inventory")
        bullets = self._extract_bullets_from_content(content) or [
            "Current technology stack assessed",
            "Data infrastructure evaluated",
            "Integration points identified",
            "Gaps and opportunities mapped",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_maturity_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add AI maturity assessment slide."""
        slide = self._add_slide_with_title(
            prs,
            "AI Maturity Assessment",
            "Current position on AI adoption curve"
        )

        content = self._extract_content_section(synthesis, "04_maturity_assessment")
        bullets = self._extract_bullets_from_content(content, max_bullets=5) or [
            "Data readiness and quality assessment",
            "Technical infrastructure evaluation",
            "Organizational AI capabilities",
            "Process automation opportunities",
            "Strategic alignment with AI initiatives",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_pain_points_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add key pain points slide."""
        slide = self._add_slide_with_title(
            prs,
            "Key Pain Points",
            "Areas where AI can deliver highest impact"
        )

        content = self._extract_content_section(synthesis, "02_pain_points")
        bullets = self._extract_bullets_from_content(content, max_bullets=6) or [
            "Manual processes requiring automation",
            "Data silos limiting insights",
            "Customer experience opportunities",
            "Operational efficiency gaps",
            "Decision-making bottlenecks",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_quick_wins_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add quick wins slide."""
        slide = self._add_slide_with_title(
            prs,
            "Top 5 Quick Wins",
            "High-impact, low-effort implementations"
        )

        content = self._extract_content_section(synthesis, "06_quick_wins")
        bullets = self._extract_bullets_from_content(content, max_bullets=5) or [
            "Deploy AI-powered customer support chatbot",
            "Implement document processing automation",
            "Launch internal knowledge base with AI search",
            "Automate routine data entry and validation",
            "Enable AI-assisted content generation",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_roadmap_overview_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add roadmap overview slide."""
        slide = self._add_slide_with_title(
            prs,
            "Implementation Roadmap",
            "Strategic timeline for AI adoption"
        )

        headers = ["Phase", "Timeline", "Focus Areas", "Expected Outcomes"]
        rows = [
            ["Foundation", "0-30 days", "Quick wins, governance", "Early value demonstration"],
            ["Build", "30-90 days", "Core capabilities", "Operational efficiency"],
            ["Scale", "90-180 days", "Expand use cases", "Revenue impact"],
            ["Transform", "180-360 days", "Enterprise AI", "Competitive advantage"],
        ]

        self._add_table_to_slide(slide, headers, rows)

    def _add_roi_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add ROI overview slide."""
        slide = self._add_slide_with_title(
            prs,
            "Expected ROI",
            "Financial impact of AI implementation"
        )

        content = self._extract_content_section(synthesis, "09_roi_calculator")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(
                slide,
                table["headers"][:4],
                table["rows"][:5],
            )
        else:
            # Fallback summary
            headers = ["Category", "Potential Savings", "Timeline"]
            rows = [
                ["Process Automation", "$50K-200K/year", "6-12 months"],
                ["Customer Experience", "15-25% improvement", "3-6 months"],
                ["Employee Productivity", "20-30% efficiency gain", "3-6 months"],
                ["Error Reduction", "40-60% decrease", "3-6 months"],
            ]
            self._add_table_to_slide(slide, headers, rows)

    def _add_next_steps_slide(self, prs: Presentation, company_input: CompanyInput) -> None:
        """Add next steps slide."""
        slide = self._add_slide_with_title(
            prs,
            "Recommended Next Steps",
            "Getting started with your AI journey"
        )

        bullets = [
            "Review and approve AI governance framework",
            "Prioritize quick wins for immediate implementation",
            "Establish AI Center of Excellence / task force",
            "Begin vendor evaluation for selected use cases",
            "Launch pilot project with measurable KPIs",
            "Schedule quarterly progress reviews",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_contact_slide(self, prs: Presentation) -> None:
        """Add contact/closing slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLORS["dark_blue"]
        p.alignment = PP_ALIGN.CENTER

        # Generated by line
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Generated by AI Strategy Factory"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS["gray"]
        p.alignment = PP_ALIGN.CENTER

    def _add_agenda_slide(self, prs: Presentation) -> None:
        """Add agenda slide for full presentation."""
        slide = self._add_slide_with_title(prs, "Agenda")

        agenda_items = [
            "Company Overview & Market Context",
            "Current State Assessment",
            "AI Maturity & Readiness Evaluation",
            "AI Use Cases & Opportunities",
            "Strategic Roadmap & Quick Wins",
            "Financial Analysis & ROI",
            "Governance & Operations Framework",
            "Change Management & Training Plan",
        ]
        self._add_bullets_to_slide(slide, agenda_items)

    # =========================================================================
    # Full findings specific slides
    # =========================================================================

    def _add_company_profile_slide(self, prs: Presentation, research: ResearchOutput) -> None:
        """Add company profile slide."""
        slide = self._add_slide_with_title(prs, "Company Profile")

        profile = research.profile
        bullets = [
            f"Description: {profile.description[:100]}..." if len(profile.description) > 100 else f"Description: {profile.description}",
            f"Business Model: {profile.business_model}" if profile.business_model else "Business Model: N/A",
            f"Size: {profile.company_size.value.title()}" if profile.company_size else "Size: Unknown",
            f"Headquarters: {profile.headquarters}" if profile.headquarters else "Headquarters: N/A",
        ]
        if profile.products_services:
            bullets.append(f"Products/Services: {', '.join(profile.products_services[:3])}")

        self._add_bullets_to_slide(slide, [b for b in bullets if b])

    def _add_industry_context_slide(self, prs: Presentation, research: ResearchOutput) -> None:
        """Add industry context slide."""
        slide = self._add_slide_with_title(prs, "Industry Context")

        industry = research.industry
        bullets = [
            f"Industry: {industry.primary_industry}" if industry.primary_industry else "Industry analysis conducted",
        ]
        if industry.key_trends:
            bullets.extend([f"Trend: {t}" for t in industry.key_trends[:4]])

        self._add_bullets_to_slide(slide, bullets)

    def _add_competitive_landscape_slide(self, prs: Presentation, research: ResearchOutput) -> None:
        """Add competitive landscape slide."""
        slide = self._add_slide_with_title(prs, "Competitive Landscape")

        if research.competitors:
            headers = ["Competitor", "AI Initiatives", "Position"]
            rows = []
            for comp in research.competitors[:5]:
                ai_init = comp.ai_initiatives[0] if comp.ai_initiatives else "Unknown"
                rows.append([comp.name, ai_init[:40], comp.market_position or "N/A"])
            self._add_table_to_slide(slide, headers, rows)
        else:
            bullets = [
                "Competitive analysis conducted",
                "Key competitors identified",
                "AI adoption benchmarks established",
                "Market positioning opportunities mapped",
            ]
            self._add_bullets_to_slide(slide, bullets)

    def _add_tech_inventory_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add technology inventory slides."""
        slide = self._add_slide_with_title(
            prs,
            "Technology Inventory",
            "Current technology landscape"
        )

        content = self._extract_content_section(synthesis, "01_tech_inventory")
        bullets = self._extract_bullets_from_content(content, max_bullets=6) or [
            "Core systems and platforms identified",
            "Data infrastructure assessed",
            "Integration capabilities evaluated",
            "Security posture reviewed",
            "Cloud readiness determined",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_pain_points_detailed_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add detailed pain points slide."""
        slide = self._add_slide_with_title(
            prs,
            "Pain Point Analysis",
            "Department-specific challenges"
        )

        content = self._extract_content_section(synthesis, "02_pain_points")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(
                slide,
                table["headers"][:4],
                table["rows"][:6],
            )
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_current_state_diagram_slide(
        self,
        prs: Presentation,
        mermaid_images: Optional[Dict[str, str]],
    ) -> None:
        """Add current state diagram slide."""
        slide = self._add_slide_with_title(prs, "Current State Architecture")

        if mermaid_images and "current_state" in mermaid_images:
            image_path = Path(mermaid_images["current_state"])
            if image_path.exists() and image_path.suffix == ".png":
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(1), Inches(1.5),
                    width=Inches(11),
                )
                return

        # Fallback message
        textbox = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(9), Inches(2)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "See mermaid_diagrams.md for architecture diagram"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["gray"]
        p.alignment = PP_ALIGN.CENTER

    def _add_maturity_detailed_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add detailed maturity assessment slides."""
        slide = self._add_slide_with_title(
            prs,
            "AI Maturity Assessment",
            "Detailed capability evaluation"
        )

        content = self._extract_content_section(synthesis, "04_maturity_assessment")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(slide, table["headers"][:5], table["rows"][:6])
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_use_case_library_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add use case library slides."""
        slide = self._add_slide_with_title(
            prs,
            "AI Use Case Library",
            "Department-specific opportunities"
        )

        content = self._extract_content_section(synthesis, "14_use_case_library")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(slide, table["headers"][:4], table["rows"][:6])
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_quick_wins_detailed_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add detailed quick wins slides."""
        slide = self._add_slide_with_title(
            prs,
            "Quick Wins - Detailed Analysis",
            "Implementation priorities"
        )

        content = self._extract_content_section(synthesis, "06_quick_wins")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(slide, table["headers"][:5], table["rows"][:6])
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_roadmap_detailed_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add detailed roadmap slides."""
        slide = self._add_slide_with_title(
            prs,
            "Implementation Roadmap - Detailed",
            "30/60/90/180/360 day plan"
        )

        content = self._extract_content_section(synthesis, "05_roadmap")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(slide, table["headers"][:4], table["rows"][:8])
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_future_state_diagram_slide(
        self,
        prs: Presentation,
        mermaid_images: Optional[Dict[str, str]],
    ) -> None:
        """Add future state diagram slide."""
        slide = self._add_slide_with_title(prs, "Future State Architecture")

        if mermaid_images and "future_state" in mermaid_images:
            image_path = Path(mermaid_images["future_state"])
            if image_path.exists() and image_path.suffix == ".png":
                slide.shapes.add_picture(
                    str(image_path),
                    Inches(1), Inches(1.5),
                    width=Inches(11),
                )
                return

        textbox = slide.shapes.add_textbox(
            Inches(2), Inches(3), Inches(9), Inches(2)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "See mermaid_diagrams.md for future state diagram"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS["gray"]
        p.alignment = PP_ALIGN.CENTER

    def _add_roi_detailed_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add detailed ROI slides."""
        slide = self._add_slide_with_title(
            prs,
            "ROI Analysis - Detailed",
            "Cost-benefit breakdown"
        )

        content = self._extract_content_section(synthesis, "09_roi_calculator")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(slide, table["headers"][:5], table["rows"][:7])
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_vendor_comparison_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add vendor comparison slides."""
        slide = self._add_slide_with_title(
            prs,
            "Vendor Comparison",
            "Build vs Buy analysis"
        )

        content = self._extract_content_section(synthesis, "07_vendor_comparison")
        table = self._extract_table_from_content(content)

        if table:
            self._add_table_to_slide(slide, table["headers"][:5], table["rows"][:6])
        else:
            bullets = self._extract_bullets_from_content(content, max_bullets=6)
            self._add_bullets_to_slide(slide, bullets)

    def _add_governance_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add governance framework slides."""
        slide = self._add_slide_with_title(
            prs,
            "AI Governance Framework",
            "Policy and compliance recommendations"
        )

        content = self._extract_content_section(synthesis, "10_ai_policy")
        bullets = self._extract_bullets_from_content(content, max_bullets=6) or [
            "AI acceptable use policy framework",
            "Data governance requirements",
            "Risk management guidelines",
            "Ethical AI principles",
            "Compliance monitoring approach",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_change_management_slides(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add change management slides."""
        slide = self._add_slide_with_title(
            prs,
            "Change Management",
            "Training and adoption strategy"
        )

        content = self._extract_content_section(synthesis, "15_change_management")
        bullets = self._extract_bullets_from_content(content, max_bullets=6) or [
            "Stakeholder communication plan",
            "Training program structure",
            "Adoption metrics and KPIs",
            "Champion network development",
            "Feedback and iteration process",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_summary_slide(self, prs: Presentation, synthesis: SynthesisOutput) -> None:
        """Add summary slide."""
        slide = self._add_slide_with_title(prs, "Summary")

        bullets = [
            "Comprehensive AI readiness assessment completed",
            "Quick wins identified for immediate value",
            "Strategic roadmap with clear milestones",
            "ROI framework with quantified benefits",
            "Governance and change management in place",
            "Ready to accelerate AI transformation",
        ]
        self._add_bullets_to_slide(slide, bullets)

    def _add_appendix_slide(self, prs: Presentation) -> None:
        """Add appendix slide."""
        slide = self._add_slide_with_title(prs, "Appendix")

        bullets = [
            "Full deliverables available in markdown format",
            "Detailed vendor comparison analysis",
            "Complete use case library by department",
            "Prompt library starter kit",
            "Glossary of AI terms",
            "Research sources and citations",
        ]
        self._add_bullets_to_slide(slide, bullets)


# Convenience functions

def generate_executive_deck(
    company_slug: str,
    company_input: CompanyInput,
    research: ResearchOutput,
    synthesis: SynthesisOutput,
    mermaid_images: Optional[Dict[str, str]] = None,
    output_dir: Optional[Path] = None,
) -> str:
    """Generate executive summary deck."""
    generator = PowerPointGenerator(output_dir=output_dir)
    return generator.generate_executive_summary(
        company_slug, company_input, research, synthesis, mermaid_images
    )


def generate_full_findings_deck(
    company_slug: str,
    company_input: CompanyInput,
    research: ResearchOutput,
    synthesis: SynthesisOutput,
    mermaid_images: Optional[Dict[str, str]] = None,
    output_dir: Optional[Path] = None,
) -> str:
    """Generate full findings presentation."""
    generator = PowerPointGenerator(output_dir=output_dir)
    return generator.generate_full_findings(
        company_slug, company_input, research, synthesis, mermaid_images
    )
